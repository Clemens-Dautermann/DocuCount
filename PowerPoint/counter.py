from pptx import Presentation
from tkinter import Tk
from tkinter.filedialog import askopenfilename

def count(p):

    #list for all strings
    text = []

    #get all text
    for slide in p.slides: #each slide
        for shape in slide.shapes:  #each shape
            if(not shape.has_text_frame): #continue if the shape has no text frames
                continue
            for paragraph in shape.text_frame.paragraphs: #for each paragraph

                for run in paragraph.runs: #for each run
                    text.append(run.text)  #get text and store

    return text


#####MAIN#####


pres = "null"

while(pres == "null"):

    #get file chooser
    Tk().withdraw()
    filename = askopenfilename()

    if(filename == ""):
        break

    try:

        print("reading file...")

        #get presentation
        pres = Presentation(filename)

        #count file
        text = count(pres)
        print("counting presentation please wait...")


        #create lists to store all words
        allWords = []
        allWordsTemp = []

        #create variables to count all chars (without spaces)
        CharsWhithoutSpacesTemp = []
        Charcount = 0

        #create variables to count chars with spaces



        #each text line
        for n in range(len(text)):

            tempLine = text[n]

            #each word in tempLine
            for i in range(len(tempLine.split(" "))):

                #add each word
                tempWordsSplit = tempLine.split(" ")
                allWordsTemp.append(tempWordsSplit[i])

        #remove word if its just a space
        for n in range(len(allWordsTemp)):

            #if the word is a space it may not be counted
            if(allWordsTemp[n] == "" or allWordsTemp[n] == " "):
                    pass

            else:
                allWords.append(allWordsTemp[n])





        #each word in the presentation without spaces
        for i in range(len(allWords)):

            #add word as charlist to temp list
            CharsWhithoutSpacesTemp.append(list(allWords[i]))

        #for each charlist in the temp list (a charlist is one word)
        for i in range(len(CharsWhithoutSpacesTemp)):

            Charcount = Charcount + len(CharsWhithoutSpacesTemp)





        print("-----------------------------------------------------------------")
        print("Total words: " + str(len(allWords)))
        print("Total lines: " + str(len(text)))
        print("Total caracters (without spaces): " + str(Charcount))


    except Exception as exc:

        print("Error! Dod you chose a .pptx file?")
        pres = "null"
        print(exc)
